"""
Generate Kippy Bootcamp Portfolio Presentation (.pptx)
Follows the required slide format:
  1. Introduction (Self-Overview, Education, Working)
  2. Overview Project (all bootcamp projects)
  3. Main Project (Kippy – background, detail, dev process, issues)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Colour palette ──────────────────────────────────────────────
GREEN  = RGBColor(0x7B, 0xC7, 0x4D)   # Kippy green
DARK   = RGBColor(0x1A, 0x1A, 0x1A)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT  = RGBColor(0xF5, 0xF5, 0xF0)   # broken white
GREY   = RGBColor(0x66, 0x66, 0x66)
LIME   = RGBColor(0x8D, 0xEE, 0x10)   # accent

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height

# ── Helper functions ────────────────────────────────────────────
def add_green_bar(slide):
    """Top green accent bar"""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.12))
    bar.fill.solid()
    bar.fill.fore_color.rgb = GREEN
    bar.line.fill.background()

def set_slide_bg(slide, color=LIGHT):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_text(slide, text, left, top, width, height, font_size=36, color=DARK, bold=True, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return tf

def add_body_text(slide, text, left, top, width, height, font_size=16, color=GREY, bold=False, alignment=PP_ALIGN.LEFT, line_spacing=1.5):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    p.space_after = Pt(6)
    return tf

def add_bullet_list(slide, items, left, top, width, height, font_size=15, color=GREY):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"•  {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(4)
    return tf

def add_card(slide, left, top, width, height, fill_color=WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_section_badge(slide, text, left, top, width=Inches(2.5), height=Inches(0.45)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = GREEN
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(2)

def add_page_number(slide, num, total):
    add_body_text(slide, f"{num} / {total}", Inches(12.2), Inches(7.0), Inches(1), Inches(0.4), font_size=10, color=GREY, alignment=PP_ALIGN.RIGHT)


TOTAL_SLIDES = 10

# ═══════════════════════════════════════════════════════════════
# SLIDE 1 – COVER
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, GREEN)

# Big title
add_title_text(slide, "KIPPY", Inches(1), Inches(1.8), Inches(11), Inches(1.4), font_size=72, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_body_text(slide, "Photo & Meme Sharing App  🐸", Inches(1), Inches(3.2), Inches(11), Inches(0.8), font_size=28, color=WHITE, bold=False, alignment=PP_ALIGN.CENTER)

# Divider line
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5), Inches(4.1), Inches(3.3), Inches(0.04))
line.fill.solid()
line.fill.fore_color.rgb = WHITE
line.line.fill.background()

# Subtitle
add_body_text(slide, "Final Project – Flutter Mobile Development Bootcamp", Inches(1), Inches(4.4), Inches(11), Inches(0.6), font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.CENTER)
add_body_text(slide, "[Nama Kamu]  •  2026", Inches(1), Inches(5.0), Inches(11), Inches(0.6), font_size=16, color=WHITE, bold=False, alignment=PP_ALIGN.CENTER)

# Logo if available
icon_path = os.path.join(os.path.dirname(__file__), "assets", "icon", "kippy_icon.png")
if os.path.exists(icon_path):
    slide.shapes.add_picture(icon_path, Inches(5.9), Inches(5.8), Inches(1.5), Inches(1.5))

# ═══════════════════════════════════════════════════════════════
# SLIDE 2 – TABLE OF CONTENTS
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_green_bar(slide)
add_title_text(slide, "Table of Contents", Inches(0.8), Inches(0.5), Inches(6), Inches(0.8), font_size=32, color=DARK)

toc_items = [
    ("01", "Introduction", "Self-overview, education & working experience"),
    ("02", "Overview Project", "Summary of all bootcamp projects"),
    ("03", "Main Project – Kippy", "Background, detail, development, issues"),
]

for i, (num, title, desc) in enumerate(toc_items):
    y = Inches(1.8) + Inches(1.6) * i
    add_card(slide, Inches(1.5), y, Inches(10), Inches(1.3))
    add_title_text(slide, num, Inches(1.9), y + Inches(0.15), Inches(1), Inches(0.7), font_size=36, color=GREEN, bold=True)
    add_title_text(slide, title, Inches(3.0), y + Inches(0.15), Inches(6), Inches(0.5), font_size=22, color=DARK, bold=True)
    add_body_text(slide, desc, Inches(3.0), y + Inches(0.7), Inches(7), Inches(0.5), font_size=14, color=GREY)

add_page_number(slide, 2, TOTAL_SLIDES)

# ═══════════════════════════════════════════════════════════════
# SLIDE 3 – INTRODUCTION: SELF-OVERVIEW
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_green_bar(slide)
add_section_badge(slide, "01  INTRODUCTION", Inches(0.8), Inches(0.5))
add_title_text(slide, "Self-Overview", Inches(0.8), Inches(1.2), Inches(6), Inches(0.8), font_size=32, color=DARK)

# Photo placeholder
photo_card = add_card(slide, Inches(1), Inches(2.3), Inches(3), Inches(3.8))
add_body_text(slide, "📷\nYour Photo Here", Inches(1.3), Inches(3.3), Inches(2.4), Inches(1.5), font_size=18, color=GREY, alignment=PP_ALIGN.CENTER)

# Info on the right
add_title_text(slide, "[Nama Lengkap Kamu]", Inches(4.8), Inches(2.3), Inches(7), Inches(0.6), font_size=26, color=DARK, bold=True)
add_body_text(slide, "Flutter Mobile Developer  •  [Kota, Indonesia]", Inches(4.8), Inches(2.9), Inches(7), Inches(0.5), font_size=16, color=GREEN, bold=True)

bio = (
    "Saya adalah seorang mahasiswa/profesional yang passionate terhadap mobile development, "
    "khususnya menggunakan Flutter & Dart. Saat ini saya mengikuti Bootcamp Flutter "
    "Mobile Development untuk mengasah kemampuan saya dalam membangun aplikasi mobile "
    "yang modern, scalable, dan memiliki UI/UX yang menarik."
)
add_body_text(slide, bio, Inches(4.8), Inches(3.6), Inches(7.5), Inches(2), font_size=14, color=GREY)

add_bullet_list(slide, [
    "Tertarik pada: Mobile Development, UI/UX Design, Clean Architecture",
    "Nilai: Kolaboratif, detail-oriented, selalu ingin belajar",
    "Tujuan karir: Menjadi Flutter developer yang handal di industri",
], Inches(4.8), Inches(4.8), Inches(7.5), Inches(2), font_size=13, color=GREY)

add_page_number(slide, 3, TOTAL_SLIDES)

# ═══════════════════════════════════════════════════════════════
# SLIDE 4 – INTRODUCTION: EDUCATION & WORKING
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_green_bar(slide)
add_section_badge(slide, "01  INTRODUCTION", Inches(0.8), Inches(0.5))

# Education section
add_title_text(slide, "🎓  Education", Inches(0.8), Inches(1.2), Inches(5), Inches(0.8), font_size=28, color=DARK)
add_card(slide, Inches(0.8), Inches(2.1), Inches(5.4), Inches(4.5))

edu_items = [
    ("[Nama Universitas / Institusi]", "[Tahun Mulai] – [Tahun Selesai]", "Jurusan: [Nama Jurusan]\nIPK: [x.xx] / 4.00\nKonsentrasi/Spesialisasi: [Bidang relevan]"),
    ("Bootcamp Flutter Mobile Dev", "2026", "Dibimbing.id\nFokusan: Flutter, Dart, Clean Architecture, Firebase, REST API"),
]

y_offset = Inches(2.3)
for title, period, desc in edu_items:
    add_title_text(slide, title, Inches(1.2), y_offset, Inches(4.5), Inches(0.4), font_size=16, color=DARK, bold=True)
    add_body_text(slide, period, Inches(1.2), y_offset + Inches(0.35), Inches(4.5), Inches(0.3), font_size=12, color=GREEN, bold=True)
    add_body_text(slide, desc, Inches(1.2), y_offset + Inches(0.65), Inches(4.5), Inches(1.2), font_size=12, color=GREY)
    y_offset += Inches(2.0)

# Working section
add_title_text(slide, "💼  Working Experience", Inches(7), Inches(1.2), Inches(5.5), Inches(0.8), font_size=28, color=DARK)
add_card(slide, Inches(7), Inches(2.1), Inches(5.5), Inches(4.5))

work_items = [
    ("[Nama Perusahaan / Organisasi]", "[Posisi]  •  [Tahun]",
     "•  Deskripsi singkat tanggung jawab\n"
     "•  Proyek yang dikerjakan\n"
     "•  Pencapaian signifikan"),
    ("[Freelance / Pengalaman Lain]", "[Posisi]  •  [Tahun]",
     "•  Deskripsi singkat tanggung jawab\n"
     "•  Tools yang digunakan"),
]

y_offset = Inches(2.3)
for title, period, desc in work_items:
    add_title_text(slide, title, Inches(7.4), y_offset, Inches(4.7), Inches(0.4), font_size=16, color=DARK, bold=True)
    add_body_text(slide, period, Inches(7.4), y_offset + Inches(0.35), Inches(4.7), Inches(0.3), font_size=12, color=GREEN, bold=True)
    add_body_text(slide, desc, Inches(7.4), y_offset + Inches(0.65), Inches(4.7), Inches(1.2), font_size=12, color=GREY)
    y_offset += Inches(2.0)

add_page_number(slide, 4, TOTAL_SLIDES)

# ═══════════════════════════════════════════════════════════════
# SLIDE 5 – OVERVIEW PROJECT
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_green_bar(slide)
add_section_badge(slide, "02  OVERVIEW PROJECT", Inches(0.8), Inches(0.5))
add_title_text(slide, "Bootcamp Projects Overview", Inches(0.8), Inches(1.2), Inches(10), Inches(0.8), font_size=32, color=DARK)

projects = [
    ("Project 1: [Nama Project]", "Deskripsi singkat project pertama. Tujuan, analysis, result, dan teknologi yang digunakan.", "Tech: Flutter, Dart"),
    ("Project 2: [Nama Project]", "Deskripsi singkat project kedua. Tujuan, analysis, result, dan teknologi yang digunakan.", "Tech: Flutter, REST API"),
    ("Project 3: [Nama Project]", "Deskripsi singkat project ketiga. Tujuan, analysis, result, dan teknologi yang digunakan.", "Tech: Flutter, Firebase"),
    ("Final Project: Kippy", "Aplikasi photo & meme sharing berbasis Flutter dengan arsitektur Clean Architecture, Firebase Auth, dan REST API integration.", "Tech: Flutter, Dart, Firebase, BLoC, Dio"),
]

for i, (title, desc, tech) in enumerate(projects):
    col = i % 2
    row = i // 2
    x = Inches(0.8) + Inches(6.1) * col
    y = Inches(2.2) + Inches(2.4) * row

    add_card(slide, x, y, Inches(5.6), Inches(2.1))
    
    # Number badge
    badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.2), y + Inches(0.2), Inches(0.5), Inches(0.5))
    badge.fill.solid()
    badge.fill.fore_color.rgb = GREEN
    badge.line.fill.background()
    tf = badge.text_frame
    p = tf.paragraphs[0]
    p.text = str(i + 1)
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    add_title_text(slide, title, x + Inches(0.9), y + Inches(0.15), Inches(4.2), Inches(0.5), font_size=17, color=DARK, bold=True)
    add_body_text(slide, desc, x + Inches(0.9), y + Inches(0.65), Inches(4.2), Inches(0.8), font_size=12, color=GREY)
    add_body_text(slide, tech, x + Inches(0.9), y + Inches(1.5), Inches(4.2), Inches(0.4), font_size=11, color=GREEN, bold=True)

add_page_number(slide, 5, TOTAL_SLIDES)

# ═══════════════════════════════════════════════════════════════
# SLIDE 6 – MAIN PROJECT: BACKGROUND & PROBLEM
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_green_bar(slide)
add_section_badge(slide, "03  MAIN PROJECT", Inches(0.8), Inches(0.5))
add_title_text(slide, "Background & Problem Statement", Inches(0.8), Inches(1.2), Inches(10), Inches(0.8), font_size=32, color=DARK)

# Left column – Background
add_card(slide, Inches(0.8), Inches(2.2), Inches(5.6), Inches(4.8))
add_title_text(slide, "🌿 Background", Inches(1.2), Inches(2.4), Inches(5), Inches(0.5), font_size=20, color=GREEN, bold=True)
bg_text = (
    "Kippy adalah aplikasi mobile berbasis Flutter untuk berbagi foto & meme "
    "dalam komunitas yang fun. Aplikasi ini dirancang sebagai Final Project "
    "Bootcamp Flutter Mobile Dev di Dibimbing.id.\n\n"
    "Target pengguna: Anak muda (Gen Z & Millennials) yang ingin berbagi "
    "konten kreatif berupa foto dan meme dalam platform yang ringan, modern, "
    "dan mudah digunakan."
)
add_body_text(slide, bg_text, Inches(1.2), Inches(3.0), Inches(4.8), Inches(3.5), font_size=14, color=GREY)

# Right column – Problem & Solution
add_card(slide, Inches(7), Inches(2.2), Inches(5.6), Inches(4.8))
add_title_text(slide, "❗ Problem", Inches(7.4), Inches(2.4), Inches(5), Inches(0.5), font_size=20, color=RGBColor(0xE0, 0x40, 0x40), bold=True)
problem_text = (
    "Platform sosial media yang ada saat ini terlalu kompleks dan overwhelming "
    "untuk sekadar berbagi foto dan meme dengan teman. Pengguna membutuhkan "
    "platform yang simpel, fokus, dan menyenangkan."
)
add_body_text(slide, problem_text, Inches(7.4), Inches(3.0), Inches(4.8), Inches(1.5), font_size=14, color=GREY)

add_title_text(slide, "✅ Solution", Inches(7.4), Inches(4.5), Inches(5), Inches(0.5), font_size=20, color=GREEN, bold=True)
solution_text = (
    "Kippy hadir sebagai solusi dengan fitur-fitur utama:\n"
    "•  Feed dinamis dengan konten foto & meme\n"
    "•  Explore & search berdasarkan genre\n"
    "•  Like, comment, dan bookmark\n"
    "•  Profil dan edit profil pengguna\n"
    "•  Activity notifications"
)
add_body_text(slide, solution_text, Inches(7.4), Inches(5.1), Inches(4.8), Inches(2), font_size=14, color=GREY)

add_page_number(slide, 6, TOTAL_SLIDES)

# ═══════════════════════════════════════════════════════════════
# SLIDE 7 – MAIN PROJECT: DETAIL APLIKASI (User Flow)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_green_bar(slide)
add_section_badge(slide, "03  MAIN PROJECT", Inches(0.8), Inches(0.5))
add_title_text(slide, "Detail Aplikasi – User Flow", Inches(0.8), Inches(1.2), Inches(10), Inches(0.8), font_size=32, color=DARK)

# Flow diagram using shapes
flow_steps = [
    ("Splash\nScreen", "Login /\nRegister", "Home\nFeed", "Explore /\nSearch", "Create\nPost", "Profile &\nSettings"),
]

step_labels = flow_steps[0]
step_w = Inches(1.6)
step_h = Inches(1.2)
start_x = Inches(0.9)
y_pos = Inches(2.5)

for i, label in enumerate(step_labels):
    x = start_x + (Inches(2.04) * i)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y_pos, step_w, step_h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = GREEN if i == 0 or i == len(step_labels)-1 else WHITE
    shape.line.color.rgb = GREEN
    shape.line.width = Pt(2)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = WHITE if i == 0 or i == len(step_labels)-1 else DARK
    p.alignment = PP_ALIGN.CENTER

    # Arrow between steps
    if i < len(step_labels) - 1:
        arrow_x = x + step_w
        arrow_w = Inches(0.44)
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, arrow_x, y_pos + Inches(0.4), arrow_w, Inches(0.4))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = GREEN
        arrow.line.fill.background()

# Screen descriptions
add_card(slide, Inches(0.8), Inches(4.2), Inches(11.7), Inches(2.8))
screen_descs = [
    ("Splash Screen", "Animasi logo Kippy saat app pertama kali dibuka"),
    ("Login / Register", "Form autentikasi user dengan email & password, social login"),
    ("Home Feed", "Scrollable feed berisi foto & meme, fitur like, comment, share"),
    ("Explore", "Grid gallery dengan filter genre dan search bar"),
    ("Create Post", "Upload foto dari gallery/camera dengan caption"),
    ("Profile", "View posts, bookmarks, edit profil, dan settings"),
]

for i, (title, desc) in enumerate(screen_descs):
    col = i % 3
    row = i // 3
    x = Inches(1.1) + Inches(3.85) * col
    y = Inches(4.4) + Inches(1.3) * row
    add_title_text(slide, f"📱 {title}", x, y, Inches(3.5), Inches(0.4), font_size=14, color=DARK, bold=True)
    add_body_text(slide, desc, x, y + Inches(0.35), Inches(3.5), Inches(0.6), font_size=12, color=GREY)

add_page_number(slide, 7, TOTAL_SLIDES)

# ═══════════════════════════════════════════════════════════════
# SLIDE 8 – MAIN PROJECT: DETAIL APLIKASI (Screenshots placeholder)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_green_bar(slide)
add_section_badge(slide, "03  MAIN PROJECT", Inches(0.8), Inches(0.5))
add_title_text(slide, "Detail Aplikasi – App Screenshots", Inches(0.8), Inches(1.2), Inches(10), Inches(0.8), font_size=32, color=DARK)
add_body_text(slide, "Tampilan utama dari berbagai layar pada aplikasi Kippy", Inches(0.8), Inches(1.9), Inches(10), Inches(0.5), font_size=16, color=GREY)

# Phone mockup placeholders
screens = ["Home Feed", "Explore", "Create Post", "Profile", "Activity"]
for i, name in enumerate(screens):
    x = Inches(0.7) + Inches(2.5) * i
    # Phone frame
    phone = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.6), Inches(2.1), Inches(3.8))
    phone.fill.solid()
    phone.fill.fore_color.rgb = WHITE
    phone.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    phone.line.width = Pt(2)
    
    add_body_text(slide, f"📱\n\n{name}\nScreenshot", x + Inches(0.2), Inches(3.5), Inches(1.7), Inches(2), font_size=14, color=GREY, alignment=PP_ALIGN.CENTER)
    add_body_text(slide, name, x, Inches(6.5), Inches(2.1), Inches(0.4), font_size=13, color=DARK, bold=True, alignment=PP_ALIGN.CENTER)

add_page_number(slide, 8, TOTAL_SLIDES)

# ═══════════════════════════════════════════════════════════════
# SLIDE 9 – MAIN PROJECT: DEVELOPMENT PROCESS & TECH STACK
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_green_bar(slide)
add_section_badge(slide, "03  MAIN PROJECT", Inches(0.8), Inches(0.5))
add_title_text(slide, "Development Process & Tech Stack", Inches(0.8), Inches(1.2), Inches(10), Inches(0.8), font_size=32, color=DARK)

# Development phases
phases = [
    ("1. Planning", "Analisis kebutuhan,\ndesain arsitektur\nClean Architecture"),
    ("2. Design", "UI/UX wireframing,\ncolor palette,\nprototyping"),
    ("3. Development", "Implementasi fitur,\nstate management\ndengan BLoC"),
    ("4. Testing", "Unit testing,\nwidget testing,\nbug fixing"),
    ("5. Polish", "Optimisasi UI,\nperformance,\nfinal review"),
]

for i, (title, desc) in enumerate(phases):
    x = Inches(0.5) + Inches(2.5) * i
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.3), Inches(2.2), Inches(2.0))
    shape.fill.solid()
    shape.fill.fore_color.rgb = GREEN if i == 2 else WHITE
    shape.line.color.rgb = GREEN
    shape.line.width = Pt(2)
    
    text_color = WHITE if i == 2 else DARK
    desc_color = WHITE if i == 2 else GREY
    add_title_text(slide, title, x + Inches(0.15), Inches(2.4), Inches(1.9), Inches(0.5), font_size=15, color=text_color, bold=True, alignment=PP_ALIGN.CENTER)
    add_body_text(slide, desc, x + Inches(0.15), Inches(2.9), Inches(1.9), Inches(1.2), font_size=12, color=desc_color, alignment=PP_ALIGN.CENTER)

# Tech Stack section
add_title_text(slide, "Tech Stack", Inches(0.8), Inches(4.7), Inches(5), Inches(0.5), font_size=22, color=DARK, bold=True)

tech_categories = [
    ("Framework", "Flutter & Dart"),
    ("State Mgmt", "BLoC (flutter_bloc)"),
    ("Architecture", "Clean Architecture"),
    ("Backend", "Firebase (Auth, Firestore, Storage)"),
    ("Networking", "Dio + REST API"),
    ("DI", "GetIt"),
    ("Local Storage", "Hive, SharedPreferences"),
    ("UI", "Lottie, CachedNetworkImage"),
]

for i, (cat, val) in enumerate(tech_categories):
    col = i % 4
    row = i // 4
    x = Inches(0.8) + Inches(3.1) * col
    y = Inches(5.3) + Inches(0.9) * row

    add_card(slide, x, y, Inches(2.8), Inches(0.75))
    add_body_text(slide, cat, x + Inches(0.15), y + Inches(0.02), Inches(2.5), Inches(0.3), font_size=11, color=GREEN, bold=True)
    add_body_text(slide, val, x + Inches(0.15), y + Inches(0.35), Inches(2.5), Inches(0.3), font_size=13, color=DARK, bold=True)

add_page_number(slide, 9, TOTAL_SLIDES)

# ═══════════════════════════════════════════════════════════════
# SLIDE 10 – MAIN PROJECT: ISSUES & PROBLEM SOLVING
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_green_bar(slide)
add_section_badge(slide, "03  MAIN PROJECT", Inches(0.8), Inches(0.5))
add_title_text(slide, "Issues & Problem Solving", Inches(0.8), Inches(1.2), Inches(10), Inches(0.8), font_size=32, color=DARK)

issues = [
    (
        "❗ API Endpoint & Authentication Error",
        "Login menghasilkan error 404 karena base URL dan endpoint yang salah di data source layer.",
        "✅ Memperbaiki konfigurasi DioClient dengan base URL yang benar, menambahkan proper error handling dan logging interceptor untuk debugging API calls."
    ),
    (
        "❗ Dark Mode Color Inconsistency",
        "Halaman Login & Register menggunakan hardcoded Colors.white/black yang tidak adapt dengan dark mode theme.",
        "✅ Mengganti semua hardcoded colors dengan theme-aware values (Theme.of(context)), kemudian disederhanakan menjadi broken white palette untuk konsistensi."
    ),
    (
        "❗ State Management untuk Feed Updates",
        "Post baru tidak langsung muncul di feed setelah di-create, harus refresh manual.",
        "✅ Mengimplementasikan ValueNotifier (globalFeedNotifier) untuk reactive state updates, sehingga new post langsung prepend ke feed list."
    ),
    (
        "❗ Navigation & Routing Complexity",
        "Integrasi Activity Page dan Edit Profile Page tidak terhubung dengan benar ke navigation flow.",
        "✅ Memperbarui AppRouter dengan routes baru, mengganti placeholder widget di MainWrapperPage dengan ActivityPage, dan menambahkan onTap callback di Settings untuk navigasi ke EditProfilePage."
    ),
]

for i, (title, problem, solution) in enumerate(issues):
    col = i % 2
    row = i // 2
    x = Inches(0.8) + Inches(6.1) * col
    y = Inches(2.1) + Inches(2.6) * row
    
    add_card(slide, x, y, Inches(5.7), Inches(2.3))
    add_title_text(slide, title, x + Inches(0.25), y + Inches(0.15), Inches(5.2), Inches(0.45), font_size=15, color=DARK, bold=True)
    add_body_text(slide, problem, x + Inches(0.25), y + Inches(0.55), Inches(5.2), Inches(0.7), font_size=12, color=GREY)
    add_body_text(slide, solution, x + Inches(0.25), y + Inches(1.4), Inches(5.2), Inches(0.8), font_size=12, color=GREEN, bold=False)

add_page_number(slide, 10, TOTAL_SLIDES)

# ═══════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════
output_path = os.path.join(os.path.dirname(__file__), "Kippy_Bootcamp_Presentation.pptx")
prs.save(output_path)
print(f"\n✅ Presentation saved to: {output_path}")
print(f"   Total slides: {TOTAL_SLIDES}")
